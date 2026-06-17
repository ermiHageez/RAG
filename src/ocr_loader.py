import io
from pathlib import Path
from typing import List, Optional
import numpy as np
from PIL import Image
from pptx import Presentation  # pyright: ignore[reportMissingImports]
from pptx.enum.shapes import MSO_SHAPE_TYPE  # pyright: ignore[reportMissingImports]
from langchain_core.documents import Document
import easyocr


_reader: Optional["easyocr.Reader"] = None


def _get_reader() -> easyocr.Reader:
    global _reader
    if _reader is None:
        _reader = easyocr.Reader(["en"])
    return _reader


def _extract_text_from_image(image_bytes: bytes) -> str:
    reader = _get_reader()
    img = Image.open(io.BytesIO(image_bytes))
    img_np = np.array(img)
    results = reader.readtext(img_np)
    return " ".join(text for _, text, _ in results)


def ocr_pdf(pdf_path: str) -> List[Document]:
    import fitz  # pyright: ignore[reportMissingImports]
    doc = fitz.open(pdf_path)
    documents = []
    for page_num in range(len(doc)):
        page = doc[page_num]
        text = page.get_text().strip()
        if not text:
            pix = page.get_pixmap(dpi=300)
            img_bytes = pix.tobytes("png")
            text = _extract_text_from_image(img_bytes)
        documents.append(Document(
            page_content=text,
            metadata={
                "source": pdf_path,
                "page": page_num + 1,
                "total_pages": len(doc),
            },
        ))
    doc.close()
    count = sum(1 for d in documents if str(d.page_content).strip())
    print(f"[DEBUG] OCR extracted {count} non-empty pages from {pdf_path}")
    return documents


def ocr_pptx(pptx_path: str) -> List[Document]:
    path = Path(pptx_path)
    prs = Presentation(str(path))
    documents = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
             if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture = shape  # pyright: ignore[reportAttributeAccessIssue]
                try:
                    text = _extract_text_from_image(picture.image.blob)  # pyright: ignore[reportAttributeAccessIssue]
                    if text.strip():
                        slide_texts.append(text)
                except Exception as e:
                    print(f"[WARN] OCR failed for slide {slide_num}: {e}")

        combined_text = "\n".join(slide_texts)
        doc = Document(
            page_content=combined_text,
            metadata={
                "source": pptx_path,
                "slide": slide_num,
                "total_slides": len(prs.slides),
            },
        )
        documents.append(doc)

    print(f"[DEBUG] OCR extracted {sum(1 for d in documents if str(d.page_content).strip())} non-empty slides from {pptx_path}")
    return documents
